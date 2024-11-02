import tkinter as tk
from tkinter import scrolledtext, filedialog, messagebox
from pptx import Presentation
from pptx.util import Inches
import os

class LyricSlideGenerator:
    def __init__(self, root):
        self.root = root
        self.root.title("Lyric to PowerPoint")
        
        # Create a label to provide instructions
        self.instruction_label = tk.Label(root, text="To split the stanza, press Enter to add a space line.")
        self.instruction_label.pack(padx=10, pady=5)

        self.lyric_text = scrolledtext.ScrolledText(root, width=50, height=10, wrap=tk.WORD)
        self.lyric_text.pack(padx=10, pady=5)

        generate_button = tk.Button(root, text="Generate PowerPoint", command=self.generate_powerpoint)
        generate_button.pack(pady=10)

        # Define template path
        self.template_path = os.path.join(os.path.expanduser("~"), "Documents", "template", "template.pptx")
        self.ensure_template_exists()

    def ensure_template_exists(self):
        # Check if the template file exists; if not, prompt the user to create it
        if not os.path.exists(self.template_path):
            response = messagebox.askyesno("Template Missing", "Template file not found. Would you like to create a template in the Documents folder?")
            if response:
                template_dir = os.path.dirname(self.template_path)
                os.makedirs(template_dir, exist_ok=True)
                
                # Create a basic template presentation
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                title = slide.shapes.title
                title.text = "Sample Title"
                
                prs.save(self.template_path)
                messagebox.showinfo("Template Created", "A template.pptx has been created in the Documents/template folder. You can now use the program.")
            else:
                messagebox.showwarning("Template Required", "Please add a template.pptx to the Documents/template folder to proceed.")
                self.root.destroy()

    def generate_powerpoint(self):
        lines = self.lyric_text.get("1.0", tk.END).strip().split("\n\n")  # Split by empty lines to identify stanzas
        prs = Presentation(self.template_path)

        for stanza in lines:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use 'Title and Content' layout

            title = slide.shapes.title
            title.text = stanza.strip()
            title.text_frame.paragraphs[0].alignment = 1  # Center alignment

            textbox = None
            for shape in slide.shapes:
                if shape.has_text_frame and shape != title:
                    textbox = shape
                    textbox.text = ""
                    textbox.text_frame.text = stanza.strip()
                    textbox.text_frame.paragraphs[0].alignment = 1  # Center alignment
                    break

            if textbox is None:
                textbox = slide.placeholders[0]
                textbox.text = ""
                textbox.text_frame.text = stanza.strip()
                textbox.width = int(Inches(9))
                textbox.height = int(Inches(5))
                textbox.left = (prs.slide_width - textbox.width) // 2
                textbox.top = (prs.slide_height - textbox.height) // 2

            textbox.vertical_anchor = 0.5  # Middle vertical alignment
            textbox.text = stanza.strip()

        file_path = self.get_unique_filename(".pptx")
        if file_path:
            prs.save(file_path)
            messagebox.showinfo("Success", f"PowerPoint saved successfully at {file_path}")

    def get_unique_filename(self, extension):
        initial_file = "Untitled" + extension
        file_path = filedialog.asksaveasfilename(initialfile=initial_file, defaultextension=extension, filetypes=[("PowerPoint files", "*" + extension)])
        if file_path:
            if os.path.exists(file_path):
                messagebox.showerror("Error", "File already exists. Please choose a different name.")
                return None
            return file_path

if __name__ == "__main__":
    root = tk.Tk()
    app = LyricSlideGenerator(root)
    root.mainloop()
